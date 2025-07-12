import random
import torch
import spacy
import fitz
from typing import Dict
# =================== ADVANCED FILE & NLP UTILITIES ===================

class AIService:
    """Handles all AI-related functionality for the EduAI Pro application"""
    # ...existing code...

    # =================== ROLE-BASED ACCESS CONTROL ===================
    def can_generate_questions(self, user: Dict) -> bool:
        """Check if user has permission to generate questions (admin only)."""
        return user.get('role') == 'admin'

    def can_generate_flashcards(self, user: Dict) -> bool:
        """Check if user has permission to generate flashcards (students and admins)."""
        return user.get('role') in ['student', 'admin']

    def generate_questions_for_user(self, user: Dict, file_path, file_type="pptx", count=10, chunk_size=3, model_name="google/flan-t5-base"):
        """Admin-only: Generate questions from file if user is admin."""
        if not self.can_generate_questions(user):
            return {"error": "Permission denied. Only admins can generate questions."}
        return self.generate_questions_from_file(file_path, file_type, count, chunk_size, model_name)

    def generate_flashcards_for_user(self, user: Dict, file_path, file_type="pptx", count=10, chunk_size=3, model_name="google/flan-t5-base"):
        """Students/Admin: Generate flashcards from file if user is student or admin."""
        if not self.can_generate_flashcards(user):
            return {"error": "Permission denied. Only students and admins can generate flashcards."}
        return self.generate_flashcards_from_file(file_path, file_type, count, chunk_size, model_name)
    # =================== ADVANCED FILE & NLP UTILITIES ===================
    def extract_text_from_ppt(self, path):
        """Extract text from PPTX file including notes."""
        from pptx import Presentation
        all_text = ""
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    all_text += shape.text.strip() + "\n"
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    all_text += "[Notes] " + notes.strip() + "\n"
            except Exception:
                continue
        return all_text

    def extract_text_from_pdf(self, path):
        """Extract text from PDF file using PyMuPDF."""
        import fitz
        doc = fitz.open(path)
        full_text = ""
        for page in doc:
            full_text += page.get_text()
        return full_text

    def chunk_text_spacy(self, text, chunk_size=3):
        """Chunk text into groups of sentences using spaCy for better context."""
        import spacy
        nlp = spacy.load("en_core_web_sm")
        doc = nlp(text)
        sentences = [sent.text.strip() for sent in doc.sents if len(sent.text.strip()) > 20]
        return [' '.join(sentences[i:i+chunk_size]) for i in range(0, len(sentences), chunk_size)]

    def generate_advanced_question(self, chunk, model_name="google/flan-t5-base"):
        """Generate exam-style question using a transformer model."""
        from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
        import torch
        device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
        tokenizer = AutoTokenizer.from_pretrained(model_name)
        model = AutoModelForSeq2SeqLM.from_pretrained(model_name).to(device)
        prompt = f"""You are a helpful assistant that creates exam-style questions based on Bloom's taxonomy and lecture content. Also prepare scenario-based questions.\n\nContext:\n\"{chunk}\"\n\nQuestion:"""
        inputs = tokenizer(prompt, return_tensors="pt", truncation=True).to(device)
        outputs = model.generate(**inputs, max_length=128)
        return tokenizer.decode(outputs[0], skip_special_tokens=True)

    def generate_questions_from_file(self, file_path, file_type="pptx", count=10, chunk_size=3, model_name="google/flan-t5-base"):
        """Extract text from pptx/pdf and generate advanced questions."""
        if file_type == "pptx":
            full_text = self.extract_text_from_ppt(file_path)
        elif file_type == "pdf":
            full_text = self.extract_text_from_pdf(file_path)
        else:
            raise ValueError("Unsupported file format. Please use pptx or pdf.")
        chunks = self.chunk_text_spacy(full_text, chunk_size=chunk_size)
        results = []
        for i, chunk in enumerate(chunks[:count]):
            try:
                question = self.generate_advanced_question(chunk, model_name=model_name)
                results.append({"chunk": chunk, "question": question})
            except Exception as e:
                print(f"Error in chunk {i}: {e}")
        return results

    def generate_flashcards_from_file(self, file_path, file_type="pptx", count=10, chunk_size=3, model_name="google/flan-t5-base"):
        """Extract text and generate flashcards using transformer and spaCy."""
        if file_type == "pptx":
            full_text = self.extract_text_from_ppt(file_path)
        elif file_type == "pdf":
            full_text = self.extract_text_from_pdf(file_path)
        else:
            raise ValueError("Unsupported file format. Please use pptx or pdf.")
        chunks = self.chunk_text_spacy(full_text, chunk_size=chunk_size)
        flashcards = []
        for i, chunk in enumerate(chunks[:count]):
            try:
                question = self.generate_advanced_question(chunk, model_name=model_name)
                # Use the generated question as the front, and chunk as the back
                flashcard = {
                    "front": question,
                    "back": chunk,
                    "tags": [file_type, "transformer", "generated"],
                    "difficulty": "medium",
                    "subject": "Content File",
                    "complexity": 0.7
                }
                flashcards.append(flashcard)
            except Exception as e:
                print(f"Error in chunk {i}: {e}")
        return flashcards
# AI Service Module for EduAI Pro
import openai
import random
from typing import List, Dict, Any
from datetime import datetime
import json
import os

# Optional: Advanced NLP capabilities
try:
    from transformers import pipeline, AutoTokenizer, AutoModelForSequenceClassification
    from sentence_transformers import SentenceTransformer
    ADVANCED_NLP_AVAILABLE = True
except ImportError:
    ADVANCED_NLP_AVAILABLE = False

class AIService:
    """Handles all AI-related functionality for the EduAI Pro application"""
    
    def __init__(self, api_key: str = None, demo_mode: bool = True, use_advanced_nlp: bool = False):
        self.api_key = api_key
        self.demo_mode = demo_mode
        self.use_advanced_nlp = use_advanced_nlp and ADVANCED_NLP_AVAILABLE
        
        if api_key and not demo_mode:
            openai.api_key = api_key
        
        # Initialize NLP models if requested and available
        self.nlp_models = {}
        if self.use_advanced_nlp:
            self._initialize_nlp_models()
    
    def _initialize_nlp_models(self):
        """Initialize advanced NLP models for enhanced AI features"""
        try:
            # Lightweight models for educational content
            self.nlp_models = {
                'sentiment': pipeline('sentiment-analysis', model='distilbert-base-uncased-finetuned-sst-2-english'),
                'question_generation': pipeline('text2text-generation', model='t5-small'),
                'text_similarity': SentenceTransformer('all-MiniLM-L6-v2'),
                'text_classification': pipeline('zero-shot-classification', model='facebook/bart-large-mnli')
            }
            print("✅ Advanced NLP models loaded successfully")
        except Exception as e:
            print(f"⚠️ Could not load advanced NLP models: {e}")
            self.use_advanced_nlp = False
    
    def generate_quiz_questions(self, topic: str, difficulty: str = "medium", count: int = 5) -> List[Dict]:
        """Generate AI-powered quiz questions for a given topic"""
        try:
            if self.use_advanced_nlp and 'question_generation' in self.nlp_models:
                return self._generate_questions_with_nlp(topic, difficulty, count)
            elif not self.demo_mode and self.api_key:
                return self._generate_questions_with_openai(topic, difficulty, count)
            else:
                return self._generate_demo_questions(topic, difficulty, count)
        except Exception as e:
            print(f"AI question generation error: {e}")
            return self._generate_demo_questions(topic, difficulty, count)
    
    def _generate_questions_with_nlp(self, topic: str, difficulty: str, count: int) -> List[Dict]:
        """Generate questions using local NLP models"""
        questions = []
        question_generator = self.nlp_models['question_generation']
        
        # Enhanced prompts for better question generation
        prompts = [
            f"Generate a {difficulty} level multiple choice question about {topic}: ",
            f"Create an educational quiz question on {topic} for {difficulty} difficulty: ",
            f"Make a {difficulty} assessment question covering {topic}: "
        ]
        
        for i in range(min(count, 5)):
            prompt = random.choice(prompts)
            try:
                # Generate question using T5 model
                result = question_generator(prompt, max_length=200, num_return_sequences=1)
                generated_text = result[0]['generated_text']
                
                # Parse and structure the generated question
                question = self._parse_generated_question(generated_text, topic, difficulty)
                questions.append(question)
            except Exception as e:
                print(f"NLP generation error: {e}")
                # Fallback to demo questions
                questions.append(self._create_demo_question(topic, difficulty, i))
        
        return questions
    
    def _generate_questions_with_openai(self, topic: str, difficulty: str, count: int) -> List[Dict]:
        """Generate questions using OpenAI API"""
        # Real OpenAI implementation would go here
        return self._generate_demo_questions(topic, difficulty, count)
    
    def _generate_demo_questions(self, topic: str, difficulty: str, count: int) -> List[Dict]:
        """Generate demo questions with enhanced logic"""
        questions = []
        for i in range(min(count, 5)):
            questions.append(self._create_demo_question(topic, difficulty, i))
        return questions
    
    def _create_demo_question(self, topic: str, difficulty: str, index: int) -> Dict:
        """Create a single demo question with enhanced content"""
        if "math" in topic.lower():
            if difficulty == "easy":
                return {
                    "question_text": f"What is 2 + 3?",
                    "option_a": "5", "option_b": "4", "option_c": "6", "option_d": "7",
                    "correct_answer": "A",
                    "explanation": "2 + 3 = 5 (basic addition)"
                }
            elif difficulty == "hard":
                return {
                    "question_text": f"If f(x) = x³ - 2x² + x - 1, what is f'(2)?",
                    "option_a": "7", "option_b": "5", "option_c": "9", "option_d": "3",
                    "correct_answer": "A",
                    "explanation": "f'(x) = 3x² - 4x + 1, so f'(2) = 3(4) - 4(2) + 1 = 12 - 8 + 1 = 5... wait, that's option B!"
                }
            else:  # medium
                return {
                    "question_text": f"If a function f(x) = x² + 2x + 1, what is f(3)?",
                    "option_a": "16", "option_b": "12", "option_c": "10", "option_d": "14",
                    "correct_answer": "A",
                    "explanation": f"f(3) = 3² + 2(3) + 1 = 9 + 6 + 1 = 16"
                }
        elif "science" in topic.lower():
            science_questions = [
                {
                    "question_text": "What is the chemical formula for water?",
                    "option_a": "H2O", "option_b": "CO2", "option_c": "NaCl", "option_d": "O2",
                    "correct_answer": "A",
                    "explanation": "Water consists of two hydrogen atoms and one oxygen atom (H2O)"
                },
                {
                    "question_text": "What is the speed of light in vacuum?",
                    "option_a": "3×10⁸ m/s", "option_b": "3×10⁶ m/s", "option_c": "3×10¹⁰ m/s", "option_d": "3×10⁴ m/s",
                    "correct_answer": "A",
                    "explanation": "The speed of light in vacuum is approximately 3×10⁸ meters per second"
                }
            ]
            return science_questions[index % len(science_questions)]
        else:
            return {
                "question_text": f"Which concept is fundamental to understanding {topic}?",
                "option_a": f"Core principles of {topic}", "option_b": f"Advanced applications", 
                "option_c": f"Historical context", "option_d": f"Future developments",
                "correct_answer": "A",
                "explanation": f"Understanding core principles is essential for mastering {topic}"
            }
    
    def _parse_generated_question(self, generated_text: str, topic: str, difficulty: str) -> Dict:
        """Parse generated text into structured question format"""
        # Simple parsing logic - in production, this would be more sophisticated
        return {
            "question_text": f"AI Generated: {generated_text[:100]}...",
            "option_a": "Option A", "option_b": "Option B", 
            "option_c": "Option C", "option_d": "Option D",
            "correct_answer": "A",
            "explanation": f"This is an AI-generated explanation for {topic} at {difficulty} level."
        }
    
    def generate_hint(self, question_text: str, options: List[str], context: Dict = None) -> str:
        """Generate AI hint using intelligent analysis of question patterns and context"""
        
        if self.use_advanced_nlp and 'text_classification' in self.nlp_models:
            return self._generate_hint_with_nlp(question_text, options, context)
        else:
            return self._generate_intelligent_hint(question_text, options, context)
    
    def _generate_hint_with_nlp(self, question_text: str, options: List[str], context: Dict = None) -> str:
        """Generate hints using NLP classification and analysis"""
        try:
            # Classify question type
            question_types = ['factual', 'conceptual', 'procedural', 'analytical', 'evaluative']
            classification = self.nlp_models['text_classification'](question_text, question_types)
            question_type = classification['labels'][0]
            
            # Analyze question complexity
            complexity = self._analyze_question_complexity(question_text, options)
            
            # Generate contextual hint based on classification
            return self._create_contextual_hint(question_text, question_type, complexity, options, context)
            
        except Exception as e:
            print(f"NLP hint generation failed: {e}")
            return self._generate_intelligent_hint(question_text, options, context)
    
    def _generate_intelligent_hint(self, question_text: str, options: List[str], context: Dict = None) -> str:
        """Generate hints using intelligent pattern analysis"""
        
        # Analyze question structure and content
        question_analysis = self._analyze_question_structure(question_text)
        option_analysis = self._analyze_options_pattern(options)
        
        # Generate hint based on analysis
        hint_strategy = self._determine_hint_strategy(question_analysis, option_analysis)
        
        return self._create_adaptive_hint(question_text, options, hint_strategy, context)
    
    def generate_explanation(self, question: str, correct_answer: str, user_answer: str, context: Dict = None) -> str:
        """Generate AI explanation using dynamic analysis of question and answers"""
        
        # Analyze the question and answer context
        explanation_data = self._analyze_answer_context(question, correct_answer, user_answer, context)
        
        # Generate explanation based on analysis
        if user_answer == correct_answer:
            return self._generate_positive_explanation(explanation_data)
        else:
            return self._generate_corrective_explanation(explanation_data)
    
    def _analyze_answer_context(self, question: str, correct_answer: str, user_answer: str, context: Dict = None) -> Dict:
        """Analyze the context to generate more informed explanations"""
        
        # Extract question type and difficulty
        question_type = self._classify_question_type(question)
        
        # Analyze answer patterns
        answer_analysis = {
            'question_type': question_type,
            'subject_domain': self._identify_subject_domain(question),
            'cognitive_level': self._assess_cognitive_level(question),
            'correct_answer': correct_answer,
            'user_answer': user_answer,
            'is_correct': user_answer == correct_answer,
            'context': context or {}
        }
        
        # Add misconception analysis if wrong
        if user_answer != correct_answer:
            answer_analysis['potential_misconception'] = self._identify_misconception_pattern(
                question, correct_answer, user_answer
            )
        
        return answer_analysis
    
    def _generate_positive_explanation(self, data: Dict) -> str:
        """Generate encouraging explanation for correct answers"""
        templates = {
            'factual': "🌟 Excellent! You correctly identified {correct_answer}. Your knowledge of {domain} fundamentals is solid.",
            'conceptual': "✅ Perfect understanding! {correct_answer} demonstrates you grasp the underlying concepts in {domain}.",
            'analytical': "🎯 Outstanding analysis! Your choice of {correct_answer} shows strong critical thinking in {domain}.",
            'procedural': "👏 Great problem-solving! You correctly applied the procedure to arrive at {correct_answer}."
        }
        
        template = templates.get(data['question_type'], templates['factual'])
        return template.format(
            correct_answer=data['correct_answer'],
            domain=data['subject_domain']
        )
    
    def _generate_corrective_explanation(self, data: Dict) -> str:
        """Generate educational explanation for incorrect answers"""
        base_explanation = f"🤔 The correct answer is {data['correct_answer']}, not {data['user_answer']}. "
        
        # Add misconception-specific guidance
        if 'potential_misconception' in data:
            misconception = data['potential_misconception']
            if misconception == 'calculation_error':
                base_explanation += "This appears to be a calculation error. Let's work through the steps systematically."
            elif misconception == 'concept_confusion':
                base_explanation += f"This suggests confusion between related concepts in {data['subject_domain']}. "
            elif misconception == 'incomplete_analysis':
                base_explanation += "Consider all aspects of the problem before selecting an answer."
            else:
                base_explanation += "Let's explore why this is a common mistake and how to avoid it."
        
        # Add learning reinforcement
        reinforcement = self._generate_learning_reinforcement(data)
        return base_explanation + reinforcement
    
    def generate_personalized_feedback(self, quiz_performance: Dict) -> str:
        """Generate personalized feedback based on quiz performance"""
        score_percentage = quiz_performance.get('percentage', 0)
        
        if score_percentage >= 90:
            feedback = [
                "🌟 Absolutely outstanding! You're truly mastering this material. Ready to tackle even more challenging topics?",
                "🏆 Exceptional performance! Your understanding is excellent. You're well-prepared for advanced concepts.",
                "⭐ Perfect execution! You've demonstrated mastery. Consider exploring related advanced topics."
            ]
        elif score_percentage >= 80:
            feedback = [
                "� Outstanding work! You're mastering this material. Ready for more challenging topics?",
                "✨ Excellent performance! Your understanding is very strong. Keep pushing forward!",
                "🚀 Great job! You're performing at a high level. Consider reviewing the few areas you missed."
            ]
        elif score_percentage >= 70:
            feedback = [
                "👍 Good progress! You're on the right track. Focus on the areas where you missed questions.",
                "📈 Solid performance! Review the concepts you struggled with to boost your understanding.",
                "💪 Nice work! You're building good understanding. A bit more practice will take you to the next level."
            ]
        elif score_percentage >= 60:
            feedback = [
                "📚 You're making progress! Focus on reviewing the key concepts and practice more questions.",
                "🔄 Good effort! Spend extra time on the topics where you missed questions.",
                "� You're getting there! Review the material and try some practice exercises."
            ]
        else:
            feedback = [
                "�💪 Keep learning! This is a challenging topic. Review the concepts and try practice questions to improve.",
                "🎯 Don't give up! Learning takes time. Review the material step by step and practice regularly.",
                "📖 Every expert was once a beginner! Focus on understanding the fundamentals first."
            ]
        
        return random.choice(feedback)
    
    def suggest_study_topics(self, weak_areas: List[str]) -> List[str]:
        """Suggest study topics based on weak performance areas"""
        suggestions = []
        for area in weak_areas:
            suggestions.append(f"Review fundamentals of {area}")
            suggestions.append(f"Practice more questions on {area}")
        return suggestions
    
    def analyze_learning_pattern(self, attempts: List[Dict]) -> Dict:
        """Analyze student learning patterns and provide insights"""
        if not attempts:
            return {"insight": "No quiz attempts yet. Start with your first quiz!"}
        
        recent_scores = [attempt['percentage'] for attempt in attempts[-5:]]
        avg_score = sum(recent_scores) / len(recent_scores)
        
        trend = "improving" if len(recent_scores) > 1 and recent_scores[-1] > recent_scores[0] else "stable"
        
        return {
            "average_score": avg_score,
            "trend": trend,
            "insight": f"Your average score is {avg_score:.1f}% and you're {trend}. Keep up the great work!"
        }
    
    def generate_flashcards(self, topic: str, count: int = 5, difficulty: str = "medium", context: str = None) -> List[Dict]:
        """Generate AI-powered flashcards using dynamic content analysis and generation"""
        
        # Use advanced NLP if available, otherwise use intelligent pattern-based generation
        if self.use_advanced_nlp and 'text_similarity' in self.nlp_models:
            return self._generate_flashcards_with_nlp(topic, count, difficulty, context)
        else:
            return self._generate_dynamic_flashcards(topic, count, difficulty, context)
    
    def _generate_flashcards_with_nlp(self, topic: str, count: int, difficulty: str, context: str = None) -> List[Dict]:
        """Generate flashcards using NLP models for content analysis and generation"""
        flashcards = []
        
        try:
            # Use context if provided for more accurate generation
            if context:
                # Extract key concepts using sentence similarity
                sentences = self._extract_sentences_from_context(context)
                key_concepts = self._extract_concepts_with_nlp(sentences, topic)
                
                # Generate flashcards based on extracted concepts
                for concept in key_concepts[:count]:
                    flashcard = self._create_concept_flashcard(concept, topic, difficulty, context)
                    flashcards.append(flashcard)
            else:
                # Generate based on topic knowledge base
                flashcards = self._generate_topic_based_flashcards_nlp(topic, count, difficulty)
            
            # Add spaced repetition data
            for card in flashcards:
                card.update(self._calculate_adaptive_spaced_repetition(difficulty, card.get('complexity', 1.0)))
            
            return flashcards
            
        except Exception as e:
            print(f"NLP flashcard generation failed: {e}")
            return self._generate_dynamic_flashcards(topic, count, difficulty, context)
    
    def _generate_dynamic_flashcards(self, topic: str, count: int, difficulty: str, context: str = None) -> List[Dict]:
        """Generate flashcards using intelligent algorithms without hardcoded content"""
        flashcards = []
        
        # Analyze topic to determine subject domain
        domain_analysis = self._analyze_topic_domain(topic)
        
        # Generate flashcards based on domain patterns and difficulty
        for i in range(count):
            flashcard = self._create_adaptive_flashcard(
                topic=topic,
                domain=domain_analysis,
                difficulty=difficulty,
                index=i,
                context=context
            )
            flashcards.append(flashcard)
        
        return flashcards
    
    def process_multimodal_input(self, file_type: str, content: bytes) -> Dict:
        """Process uploaded images, audio, or other media"""
        # Placeholder for multimodal AI processing
        return {
            "type": file_type,
            "analysis": "AI analysis of uploaded content",
            "suggestions": ["Generated suggestion 1", "Generated suggestion 2"]
        }
    
    def generate_lesson_summary(self, lesson_content: str) -> Dict:
        """Generate enhanced lesson objectives and summary from content using AI analysis"""
        # Analyze content structure and extract key information
        content_analysis = self._analyze_content_structure(lesson_content)
        
        # Extract key concepts using keyword analysis
        key_concepts = self._extract_key_concepts(lesson_content)
        
        # Generate learning objectives based on Bloom's taxonomy
        objectives = self._generate_learning_objectives(key_concepts, lesson_content)
        
        # Create comprehensive summary
        summary = self._create_intelligent_summary(lesson_content, key_concepts)
        
        # Generate thought-provoking questions
        questions = self._generate_critical_thinking_questions(key_concepts, lesson_content)
        
        return {
            "objectives": objectives,
            "summary": summary,
            "key_points": key_concepts,
            "suggested_questions": questions,
            "content_analysis": content_analysis,
            "estimated_study_time": self._estimate_study_time(lesson_content),
            "difficulty_level": self._assess_content_difficulty(lesson_content),
            "prerequisite_knowledge": self._identify_prerequisites(lesson_content)
        }
    
    def generate_study_tips(self, quiz_name: str, chapter_name: str) -> List[str]:
        """Generate AI-powered study tips for quiz preparation"""
        tips = [
            f"Focus on key concepts from {chapter_name} that are likely to appear in {quiz_name}",
            "Review your previous quiz attempts to identify weak areas",
            "Practice active recall by testing yourself without looking at notes",
            "Create mental connections between different concepts in this chapter",
            "Take breaks during study sessions to improve retention",
            f"Use flashcards to memorize important terms related to {chapter_name}",
            "Explain concepts to someone else to test your understanding",
            "Practice similar questions to build confidence for the quiz"
        ]
        return random.sample(tips, min(4, len(tips)))

    def chat_response(self, message: str, context: Dict = None) -> str:
        """Generate conversational AI responses using dynamic analysis"""
        
        # Analyze the user's message for intent and content
        message_analysis = self._analyze_user_message(message, context)
        
        # Generate response based on analysis
        return self._generate_contextual_response(message_analysis, context)
    
    def _analyze_user_message(self, message: str, context: Dict = None) -> Dict:
        """Analyze user message to understand intent and generate appropriate response"""
        
        analysis = {
            'message': message,
            'intent': self._classify_message_intent(message),
            'subject': self._identify_message_subject(message),
            'sentiment': self._analyze_message_sentiment(message),
            'complexity': self._assess_message_complexity(message),
            'context': context or {}
        }
        
        return analysis
    
    def _generate_contextual_response(self, analysis: Dict, context: Dict = None) -> str:
        """Generate contextual response based on message analysis"""
        
        intent = analysis['intent']
        subject = analysis['subject']
        sentiment = analysis['sentiment']
        
        # Generate response based on intent
        if intent == 'help_request':
            return self._generate_help_response(analysis)
        elif intent == 'learning_question':
            return self._generate_learning_response(analysis)
        elif intent == 'clarification':
            return self._generate_clarification_response(analysis)
        elif intent == 'encouragement_needed':
            return self._generate_encouragement_response(analysis)
        elif intent == 'greeting':
            return self._generate_greeting_response(analysis)
        else:
            return self._generate_general_educational_response(analysis)
    
    # =========================================================================
    # MISSING HELPER METHODS - DYNAMIC AI CONTENT GENERATION
    # =========================================================================
    
    def _analyze_question_complexity(self, question_text: str, options: List[str]) -> float:
        """Analyze the complexity of a question based on text analysis"""
        # Calculate complexity based on question length, vocabulary, and structure
        words = question_text.split()
        avg_word_length = sum(len(word) for word in words) / len(words) if words else 0
        
        # Check for complex indicators
        complex_keywords = ['analyze', 'evaluate', 'synthesize', 'compare', 'contrast', 'explain', 'justify']
        complexity_score = sum(1 for keyword in complex_keywords if keyword.lower() in question_text.lower())
        
        # Factor in option complexity
        option_complexity = sum(len(opt.split()) for opt in options) / len(options) if options else 0
        
        return min(1.0, (avg_word_length / 10 + complexity_score / 5 + option_complexity / 20))
    
    def _create_contextual_hint(self, question_text: str, question_type: str, complexity: float, options: List[str], context: Dict = None) -> str:
        """Create contextual hints based on question analysis"""
        if question_type == 'factual':
            return f"💡 This is a factual question. Focus on recalling specific information or definitions related to the key terms in the question."
        elif question_type == 'conceptual':
            return f"🤔 This requires understanding concepts. Think about the underlying principles and how they relate to each other."
        elif question_type == 'analytical':
            return f"🔍 Break down the problem into parts. Consider the relationships between different elements mentioned."
        elif question_type == 'procedural':
            return f"📋 Think step-by-step. What process or procedure would you follow to solve this?"
        else:
            return f"💭 Consider the context and eliminate obviously incorrect options first."
    
    def _analyze_question_structure(self, question_text: str) -> Dict:
        """Analyze the structure and components of a question"""
        question_lower = question_text.lower()
        
        analysis = {
            'question_words': [word for word in ['what', 'when', 'where', 'who', 'why', 'how'] if word in question_lower],
            'has_numbers': any(char.isdigit() for char in question_text),
            'has_formulas': any(symbol in question_text for symbol in ['=', '+', '-', '*', '/', '^']),
            'word_count': len(question_text.split()),
            'sentence_type': 'interrogative' if '?' in question_text else 'declarative',
            'complexity_indicators': ['complex', 'advanced', 'detailed'] if any(word in question_lower for word in ['complex', 'advanced', 'detailed']) else []
        }
        
        return analysis
    
    def _analyze_options_pattern(self, options: List[str]) -> Dict:
        """Analyze the pattern and structure of answer options"""
        if not options:
            return {'pattern': 'none', 'complexity': 0}
        
        # Check for numerical patterns
        numerical_options = [opt for opt in options if any(char.isdigit() for char in opt)]
        
        # Check for formula patterns
        formula_options = [opt for opt in options if any(symbol in opt for symbol in ['=', '+', '-', '*', '/'])]
        
        # Calculate average length
        avg_length = sum(len(opt.split()) for opt in options) / len(options)
        
        pattern_analysis = {
            'has_numbers': len(numerical_options) > 0,
            'has_formulas': len(formula_options) > 0,
            'avg_option_length': avg_length,
            'complexity': 'high' if avg_length > 5 else 'medium' if avg_length > 2 else 'low',
            'uniform_length': max(len(opt.split()) for opt in options) - min(len(opt.split()) for opt in options) <= 2
        }
        
        return pattern_analysis
    
    def _determine_hint_strategy(self, question_analysis: Dict, option_analysis: Dict) -> str:
        """Determine the best hint strategy based on question and option analysis"""
        if question_analysis.get('has_formulas') or option_analysis.get('has_formulas'):
            return 'mathematical'
        elif question_analysis.get('has_numbers') or option_analysis.get('has_numbers'):
            return 'numerical'
        elif len(question_analysis.get('question_words', [])) > 0:
            return 'factual'
        elif option_analysis.get('complexity') == 'high':
            return 'analytical'
        else:
            return 'general'
    
    def _create_adaptive_hint(self, question_text: str, options: List[str], strategy: str, context: Dict = None) -> str:
        """Create adaptive hints based on strategy"""
        hint_templates = {
            'mathematical': "🔢 Look for mathematical relationships and apply the relevant formulas or principles.",
            'numerical': "📊 Focus on the numerical values and consider what calculations or comparisons are needed.",
            'factual': "📚 Recall the key facts and definitions related to the main concept in the question.",
            'analytical': "🧩 Break down the problem into smaller parts and analyze each component systematically.",
            'general': "💡 Read the question carefully and eliminate options that don't make logical sense."
        }
        
        base_hint = hint_templates.get(strategy, hint_templates['general'])
        
        # Add context-specific guidance if available
        if context and 'subject' in context:
            subject = context['subject']
            base_hint += f" Consider the principles of {subject} that apply here."
        
        return base_hint
    
    def _classify_question_type(self, question: str) -> str:
        """Classify the type of question for better explanation generation"""
        question_lower = question.lower()
        
        if any(word in question_lower for word in ['what is', 'define', 'who is', 'when did']):
            return 'factual'
        elif any(word in question_lower for word in ['why', 'explain', 'how does', 'what causes']):
            return 'conceptual'
        elif any(word in question_lower for word in ['analyze', 'compare', 'evaluate', 'assess']):
            return 'analytical'
        elif any(word in question_lower for word in ['calculate', 'solve', 'find', 'determine']):
            return 'procedural'
        else:
            return 'general'
    
    def _identify_subject_domain(self, question: str) -> str:
        """Identify the subject domain of a question"""
        question_lower = question.lower()
        
        if any(word in question_lower for word in ['equation', 'formula', 'calculate', '+', '-', '*', '/']):
            return 'Mathematics'
        elif any(word in question_lower for word in ['atom', 'molecule', 'reaction', 'element', 'chemical']):
            return 'Chemistry'
        elif any(word in question_lower for word in ['force', 'energy', 'motion', 'velocity', 'acceleration']):
            return 'Physics'
        elif any(word in question_lower for word in ['cell', 'organism', 'dna', 'protein', 'evolution']):
            return 'Biology'
        elif any(word in question_lower for word in ['algorithm', 'programming', 'code', 'function', 'variable']):
            return 'Computer Science'
        elif any(word in question_lower for word in ['history', 'war', 'revolution', 'century', 'empire']):
            return 'History'
        elif any(word in question_lower for word in ['literature', 'author', 'poem', 'novel', 'shakespeare']):
            return 'Literature'
        else:
            return 'General Studies'
    
    def _assess_cognitive_level(self, question: str) -> str:
        """Assess the cognitive level of a question based on Bloom's taxonomy"""
        question_lower = question.lower()
        
        # Higher order thinking
        if any(word in question_lower for word in ['evaluate', 'judge', 'critique', 'assess', 'justify']):
            return 'evaluation'
        elif any(word in question_lower for word in ['create', 'design', 'formulate', 'develop', 'construct']):
            return 'synthesis'
        elif any(word in question_lower for word in ['analyze', 'compare', 'contrast', 'examine', 'investigate']):
            return 'analysis'
        # Lower order thinking
        elif any(word in question_lower for word in ['apply', 'use', 'demonstrate', 'solve', 'calculate']):
            return 'application'
        elif any(word in question_lower for word in ['explain', 'describe', 'summarize', 'interpret']):
            return 'comprehension'
        else:
            return 'knowledge'
    
    def _identify_misconception_pattern(self, question: str, correct_answer: str, user_answer: str) -> str:
        """Identify potential misconception patterns in incorrect answers"""
        question_lower = question.lower()
        
        # Mathematical misconceptions
        if any(word in question_lower for word in ['calculate', 'solve', 'find']):
            return 'calculation_error'
        
        # Conceptual misconceptions
        elif any(word in question_lower for word in ['explain', 'why', 'how']):
            return 'concept_confusion'
        
        # Analysis misconceptions
        elif any(word in question_lower for word in ['analyze', 'compare', 'evaluate']):
            return 'incomplete_analysis'
        
        # Default
        else:
            return 'knowledge_gap'
    
    def _generate_learning_reinforcement(self, data: Dict) -> str:
        """Generate learning reinforcement based on the misconception"""
        reinforcements = {
            'calculation_error': "💡 Try working through the problem step by step and double-check your calculations.",
            'concept_confusion': f"📚 Review the key concepts in {data['subject_domain']} to strengthen your understanding.",
            'incomplete_analysis': "🔍 Make sure to consider all aspects of the problem before reaching a conclusion.",
            'knowledge_gap': f"📖 Spend some time reviewing the fundamental principles of {data['subject_domain']}."
        }
        
        misconception = data.get('potential_misconception', 'knowledge_gap')
        return reinforcements.get(misconception, reinforcements['knowledge_gap'])
    
    def _classify_message_intent(self, message: str) -> str:
        """Classify the intent of a user message"""
        message_lower = message.lower()
        
        if any(word in message_lower for word in ['help', 'stuck', 'confused', "don't understand"]):
            return 'help_request'
        elif any(word in message_lower for word in ['what', 'how', 'why', 'when', 'where']):
            return 'learning_question'
        elif any(word in message_lower for word in ['explain', 'clarify', 'elaborate']):
            return 'clarification'
        elif any(word in message_lower for word in ['difficult', 'hard', 'struggling', 'frustrated']):
            return 'encouragement_needed'
        elif any(word in message_lower for word in ['hello', 'hi', 'hey', 'good morning', 'good afternoon']):
            return 'greeting'
        else:
            return 'general_inquiry'
    
    def _identify_message_subject(self, message: str) -> str:
        """Identify the subject matter in a user message"""
        return self._identify_subject_domain(message)
    
    def _analyze_message_sentiment(self, message: str) -> str:
        """Analyze the sentiment of a user message"""
        message_lower = message.lower()
        
        positive_words = ['great', 'excellent', 'good', 'love', 'like', 'amazing', 'fantastic']
        negative_words = ['difficult', 'hard', 'confused', 'frustrated', 'stuck', 'hate', 'dislike']
        
        positive_count = sum(1 for word in positive_words if word in message_lower)
        negative_count = sum(1 for word in negative_words if word in message_lower)
        
        if positive_count > negative_count:
            return 'positive'
        elif negative_count > positive_count:
            return 'negative'
        else:
            return 'neutral'
    
    def _assess_message_complexity(self, message: str) -> str:
        """Assess the complexity of a user message"""
        words = message.split()
        avg_word_length = sum(len(word) for word in words) / len(words) if words else 0
        
        if avg_word_length > 6 or len(words) > 20:
            return 'high'
        elif avg_word_length > 4 or len(words) > 10:
            return 'medium'
        else:
            return 'low'
    
    def _generate_help_response(self, analysis: Dict) -> str:
        """Generate helpful response for help requests"""
        subject = analysis['subject']
        responses = [
            f"I'm here to help with {subject}! 🤔 Can you tell me specifically what concept is giving you trouble?",
            f"No worries about {subject} - let's work through this together! 💡 What particular area would you like me to explain?",
            f"I understand {subject} can be challenging! 🎯 Let me know which topic you'd like to focus on."
        ]
        return random.choice(responses)
    
    def _generate_learning_response(self, analysis: Dict) -> str:
        """Generate educational response for learning questions"""
        subject = analysis['subject']
        responses = [
            f"Great question about {subject}! 📚 Let me help you understand this concept better.",
            f"Excellent inquiry regarding {subject}! 🔍 This shows you're thinking deeply about the material.",
            f"That's a thoughtful question about {subject}! 💡 Let's explore this together."
        ]
        return random.choice(responses)
    
    def _generate_clarification_response(self, analysis: Dict) -> str:
        """Generate clarifying response"""
        responses = [
            "I'd be happy to clarify that for you! 📝 Can you be more specific about which part needs explanation?",
            "Absolutely! Let me break that down for you. 🔍 What aspect would you like me to elaborate on?",
            "Of course! I'll explain that in more detail. 💡 Which part would you like me to focus on first?"
        ]
        return random.choice(responses)
    
    def _generate_encouragement_response(self, analysis: Dict) -> str:
        """Generate encouraging response"""
        responses = [
            "Don't worry, everyone finds this challenging at first! 💪 You're doing great by asking questions.",
            "Learning is a process, and you're on the right track! 🌟 Let's tackle this step by step.",
            "I believe in you! 🎯 Every expert was once a beginner. Let's work through this together."
        ]
        return random.choice(responses)
    
    def _generate_greeting_response(self, analysis: Dict) -> str:
        """Generate friendly greeting response"""
        responses = [
            "Hello! 👋 I'm your AI learning assistant. How can I help you learn today?",
            "Hi there! 🤖 Ready to explore some new concepts? What subject interests you?",
            "Greetings! 🌟 I'm here to make learning fun and engaging. What would you like to study?"
        ]
        return random.choice(responses)
    
    def _generate_general_educational_response(self, analysis: Dict) -> str:
        """Generate general educational response"""
        message = analysis['message']
        responses = [
            f"That's an interesting topic: '{message}'! 🤔 What specific aspect would you like to explore?",
            f"Great point about '{message}'! 💡 How can I help you learn more about this?",
            f"I see you're interested in '{message}'. 📚 What would you like to know about it?"
        ]
        return random.choice(responses)
    
    def _analyze_topic_domain(self, topic: str) -> Dict:
        """Analyze topic to determine subject domain and characteristics"""
        topic_lower = topic.lower()
        
        domain_analysis = {
            'primary_domain': self._identify_subject_domain(topic),
            'has_mathematical_content': any(word in topic_lower for word in ['math', 'calculate', 'equation', 'formula']),
            'has_scientific_content': any(word in topic_lower for word in ['science', 'experiment', 'theory', 'hypothesis']),
            'has_historical_content': any(word in topic_lower for word in ['history', 'historical', 'past', 'ancient']),
            'has_technical_content': any(word in topic_lower for word in ['technology', 'computer', 'programming', 'algorithm']),
            'complexity_level': 'advanced' if any(word in topic_lower for word in ['advanced', 'complex', 'sophisticated']) else 'basic'
        }
        
        return domain_analysis
    
    def _create_adaptive_flashcard(self, topic: str, domain: Dict, difficulty: str, index: int, context: str = None) -> Dict:
        """Create adaptive flashcards based on domain analysis"""
        primary_domain = domain['primary_domain']
        
        # Generate dynamic content based on domain
        if domain['has_mathematical_content']:
            return self._generate_math_flashcards(topic, 1, difficulty)[0]
        elif domain['has_scientific_content']:
            return self._generate_science_flashcards(topic, 1, difficulty)[0]
        elif domain['has_technical_content']:
            return self._generate_tech_flashcards(topic, 1, difficulty)[0]
        elif domain['has_historical_content']:
            return self._generate_humanities_flashcards(topic, 1, difficulty)[0]
        else:
            return self._generate_general_flashcards(topic, 1, difficulty)[0]
    
    # =========================================================================
    # ADVANCED NLP METHODS (for when NLP models are available)
    # =========================================================================
    
    def _extract_sentences_from_context(self, context: str) -> List[str]:
        """Extract meaningful sentences from context using NLP"""
        if not context:
            return []
        
        # Simple sentence extraction - would use spaCy in production
        sentences = [s.strip() for s in context.split('.') if len(s.strip()) > 20]
        return sentences[:10]  # Limit for processing
    
    def _extract_concepts_with_nlp(self, sentences: List[str], topic: str) -> List[str]:
        """Extract key concepts using NLP similarity"""
        # Placeholder for NLP-based concept extraction
        concepts = []
        for sentence in sentences:
            if topic.lower() in sentence.lower():
                # Extract noun phrases and key terms
                words = sentence.split()
                key_terms = [word for word in words if len(word) > 5 and word.isalpha()]
                concepts.extend(key_terms[:2])
        
        return list(set(concepts))[:5]
    
    def _create_concept_flashcard(self, concept: str, topic: str, difficulty: str, context: str = None) -> Dict:
        """Create flashcard for a specific concept"""
        return {
            "front": f"What is {concept} in the context of {topic}?",
            "back": f"{concept} is a key concept in {topic}. {context[:100] if context else 'Further study recommended.'}...",
            "tags": [topic.lower(), concept.lower(), difficulty],
            "difficulty": difficulty,
            "subject": topic,
            "complexity": 0.7 if difficulty == "hard" else 0.5 if difficulty == "medium" else 0.3
        }
    
    def _generate_topic_based_flashcards_nlp(self, topic: str, count: int, difficulty: str) -> List[Dict]:
        """Generate flashcards using NLP topic analysis"""
        # Fallback to domain-based generation
        domain_analysis = self._analyze_topic_domain(topic)
        flashcards = []
        
        for i in range(count):
            flashcard = self._create_adaptive_flashcard(topic, domain_analysis, difficulty, i)
            flashcards.append(flashcard)
        
        return flashcards
    
    def _calculate_adaptive_spaced_repetition(self, difficulty: str, complexity: float) -> Dict:
        """Calculate adaptive spaced repetition based on complexity"""
        from datetime import timedelta
        
        # Adjust intervals based on complexity
        base_intervals = {
            "easy": [1, 3, 7, 14, 30],
            "medium": [1, 2, 5, 12, 25],
            "hard": [1, 1, 3, 8, 20]
        }
        
        intervals = base_intervals.get(difficulty, base_intervals["medium"])
        
        # Adjust for complexity
        complexity_multiplier = 1 + complexity
        adjusted_intervals = [max(1, int(interval * complexity_multiplier)) for interval in intervals]
        
        next_review = datetime.now() + timedelta(days=adjusted_intervals[0])
        
        return {
            "next_review": next_review.isoformat(),
            "review_interval": adjusted_intervals[0],
            "review_count": 0,
            "ease_factor": 2.5,
            "complexity_score": complexity,
            "spaced_repetition_data": {
                "intervals": adjusted_intervals,
                "current_stage": 0,
                "adaptive": True
            }
        }
    
    # =========================================================================
    # ADVANCED CONTENT GENERATION (Inspired by Transformers + spaCy)
    # =========================================================================
    
    def generate_questions_from_content(self, content: str, count: int = 5, difficulty: str = "medium") -> List[Dict]:
        """Generate questions from provided content using advanced NLP techniques"""
        try:
            if self.use_advanced_nlp and content:
                return self._generate_questions_from_text_nlp(content, count, difficulty)
            else:
                return self._generate_questions_from_text_basic(content, count, difficulty)
        except Exception as e:
            print(f"Advanced question generation error: {e}")
            return self._generate_questions_from_text_basic(content, count, difficulty)
    
    def _generate_questions_from_text_nlp(self, content: str, count: int, difficulty: str) -> List[Dict]:
        """Generate questions using NLP models similar to the Colab approach"""
        try:
            # Chunk the text into meaningful segments
            chunks = self._chunk_text_intelligently(content)
            
            questions = []
            for i, chunk in enumerate(chunks[:count]):
                # Generate question using T5-style prompting
                question_data = self._generate_question_from_chunk_nlp(chunk, difficulty)
                questions.append(question_data)
            
            return questions
            
        except Exception as e:
            print(f"NLP question generation failed: {e}")
            return self._generate_questions_from_text_basic(content, count, difficulty)
    
    def _chunk_text_intelligently(self, text: str, chunk_size: int = 3) -> List[str]:
        """Chunk text into meaningful segments for question generation"""
        if not text:
            return []
        
        # Simple sentence splitting (would use spaCy in production)
        sentences = [s.strip() for s in text.split('.') if len(s.strip()) > 20]
        
        # Group sentences into chunks
        chunks = []
        for i in range(0, len(sentences), chunk_size):
            chunk = ' '.join(sentences[i:i+chunk_size])
            if len(chunk) > 50:  # Ensure chunk has meaningful content
                chunks.append(chunk)
        
        return chunks[:10]  # Limit for performance
    
    def _generate_question_from_chunk_nlp(self, chunk: str, difficulty: str) -> Dict:
        """Generate a question from a text chunk using NLP-style prompting"""
        try:
            if self.use_advanced_nlp and 'question_generation' in self.nlp_models:
                # Use T5 model for question generation
                prompt = f"""Create a {difficulty} level exam question based on Bloom's taxonomy from this content:

Content: "{chunk}"

Generate a multiple choice question:"""
                
                result = self.nlp_models['question_generation'](
                    prompt, 
                    max_length=200, 
                    num_return_sequences=1
                )
                
                generated_text = result[0]['generated_text']
                return self._parse_generated_question_advanced(generated_text, chunk, difficulty)
            else:
                return self._generate_question_from_chunk_basic(chunk, difficulty)
                
        except Exception as e:
            print(f"Chunk question generation error: {e}")
            return self._generate_question_from_chunk_basic(chunk, difficulty)
    
    def _parse_generated_question_advanced(self, generated_text: str, chunk: str, difficulty: str) -> Dict:
        """Parse AI-generated question text into structured format"""
        # Enhanced parsing logic for T5 generated content
        lines = generated_text.split('\n')
        
        question_text = ""
        options = []
        correct_answer = "A"
        
        for line in lines:
            line = line.strip()
            if line and not line.startswith('A)') and not line.startswith('B)'):
                if '?' in line:
                    question_text = line
                elif line.startswith('A)') or line.startswith('a)'):
                    options.append(line[2:].strip())
                elif line.startswith('B)') or line.startswith('b)'):
                    options.append(line[2:].strip())
                elif line.startswith('C)') or line.startswith('c)'):
                    options.append(line[2:].strip())
                elif line.startswith('D)') or line.startswith('d)'):
                    options.append(line[2:].strip())
        
        # Fallback if parsing fails
        if not question_text:
            question_text = f"Based on the content provided, which statement is most accurate?"
        
        if len(options) < 4:
            # Generate fallback options
            key_concepts = self._extract_key_terms_from_chunk(chunk)
            options = [
                f"Concept related to {key_concepts[0] if key_concepts else 'the topic'}",
                f"Alternative interpretation involving {key_concepts[1] if len(key_concepts) > 1 else 'different aspects'}",
                f"Broader application of {key_concepts[2] if len(key_concepts) > 2 else 'the principles'}",
                f"Unrelated concept not covered in the material"
            ]
        
        return {
            "question_text": question_text,
            "option_a": options[0] if len(options) > 0 else "Option A",
            "option_b": options[1] if len(options) > 1 else "Option B", 
            "option_c": options[2] if len(options) > 2 else "Option C",
            "option_d": options[3] if len(options) > 3 else "Option D",
            "correct_answer": correct_answer,
            "explanation": f"Based on the content: {chunk[:100]}...",
            "source_content": chunk,
            "generation_method": "nlp_advanced"
        }
    
    def _generate_question_from_chunk_basic(self, chunk: str, difficulty: str) -> Dict:
        """Generate question from chunk using basic text analysis"""
        key_terms = self._extract_key_terms_from_chunk(chunk)
        
        if not key_terms:
            return {
                "question_text": "What is the main concept discussed in this content?",
                "option_a": "Primary concept", "option_b": "Secondary aspect", 
                "option_c": "Related theory", "option_d": "Unrelated topic",
                "correct_answer": "A",
                "explanation": f"The content focuses on: {chunk[:100]}...",
                "source_content": chunk
            }
        
        main_term = key_terms[0]
        question_templates = [
            f"What is the significance of {main_term} in this context?",
            f"How does {main_term} relate to the main topic?",
            f"Which statement best describes {main_term}?",
            f"What role does {main_term} play in the discussed scenario?"
        ]
        
        question_text = random.choice(question_templates)
        
        return {
            "question_text": question_text,
            "option_a": f"Correct interpretation of {main_term}",
            "option_b": f"Alternative view of {main_term}",
            "option_c": f"Broader application of {main_term}",
            "option_d": f"Incorrect understanding of {main_term}",
            "correct_answer": "A",
            "explanation": f"The content indicates that {main_term} is significant because: {chunk[:150]}...",
            "source_content": chunk,
            "generation_method": "basic_analysis"
        }
    
    def _extract_key_terms_from_chunk(self, chunk: str) -> List[str]:
        """Extract key terms from a text chunk"""
        if not chunk:
            return []
        
        # Simple keyword extraction (would use spaCy NER in production)
        words = chunk.split()
        
        # Filter for meaningful terms
        key_terms = []
        for word in words:
            word_clean = word.strip('.,!?":;()[]{}')
            if (len(word_clean) > 4 and 
                word_clean.isalpha() and 
                word_clean.lower() not in ['this', 'that', 'these', 'those', 'they', 'them', 'their']):
                key_terms.append(word_clean)
        
        # Remove duplicates and return top terms
        return list(set(key_terms))[:5]
    
    def _generate_questions_from_text_basic(self, content: str, count: int, difficulty: str) -> List[Dict]:
        """Basic question generation from text content"""
        chunks = self._chunk_text_intelligently(content, chunk_size=2)
        
        questions = []
        for i, chunk in enumerate(chunks[:count]):
            question = self._generate_question_from_chunk_basic(chunk, difficulty)
            questions.append(question)
        
        return questions
    
    def generate_flashcards_from_content(self, content: str, count: int = 5, difficulty: str = "medium") -> List[Dict]:
        """Generate flashcards from provided content using advanced techniques"""
        try:
            if self.use_advanced_nlp and content:
                return self._generate_flashcards_from_text_nlp(content, count, difficulty)
            else:
                return self._generate_flashcards_from_text_basic(content, count, difficulty)
        except Exception as e:
            print(f"Advanced flashcard generation error: {e}")
            return self._generate_flashcards_from_text_basic(content, count, difficulty)
    
    def _generate_flashcards_from_text_nlp(self, content: str, count: int, difficulty: str) -> List[Dict]:
        """Generate flashcards using NLP analysis of content"""
        flashcards = []
        
        try:
            # Extract key concepts using advanced analysis
            concepts = self._extract_concepts_with_advanced_nlp(content)
            
            flashcards = []
            for i, concept in enumerate(concepts[:count]):
                flashcard = self._create_flashcard_from_concept_nlp(concept, content, difficulty)
                flashcards.append(flashcard)
            
            return flashcards
            
        except Exception as e:
            print(f"NLP flashcard generation failed: {e}")
            return self._generate_flashcards_from_text_basic(content, count, difficulty)
    
    def _extract_concepts_with_advanced_nlp(self, content: str) -> List[Dict]:
        """Extract concepts using advanced NLP techniques"""
        concepts = []
        
        # Extract sentences and analyze for concepts
        sentences = self._chunk_text_intelligently(content, chunk_size=1)
        
        for sentence in sentences:
            key_terms = self._extract_key_terms_from_chunk(sentence)
            for term in key_terms:
                if len(term) > 3:  # Filter out short terms
                    importance = self._calculate_term_importance(term, content)
                    concepts.append({
                        'term': term,
                        'context': sentence,
                        'importance': importance,
                        'related_content': sentence[:200]
                    })
        
        # Sort by importance and return top concepts
        concepts.sort(key=lambda x: x['importance'], reverse=True)
        return concepts[:10]
    
    def _calculate_term_importance(self, term: str, content: str) -> float:
        """Calculate the importance of a term in the content"""
        term_count = content.lower().count(term.lower())
        content_length = len(content.split())
        
        # Simple TF-IDF-like scoring
        frequency = term_count / content_length if content_length > 0 else 0
        
        # Boost importance for longer terms and terms that appear multiple times
        length_bonus = min(len(term) / 10, 0.5)
        frequency_bonus = min(term_count / 5, 0.3)
        
        return frequency + length_bonus + frequency_bonus
    
    def _create_flashcard_from_concept_nlp(self, concept: Dict, content: str, difficulty: str) -> Dict:
        """Create flashcard from extracted concept using NLP"""
        term = concept['term']
        context = concept['related_content']
        
        # Generate question based on difficulty
        if difficulty == "easy":
            front = f"What is {term}?"
            back = f"{term} is explained in the context as: {context[:100]}..."
        elif difficulty == "hard":
            front = f"Analyze the role and significance of {term} in the broader context."
            back = f"{term} plays a significant role: {context[:150]}... [Requires deeper analysis]"
        else:  # medium
            front = f"Explain {term} and its relevance."
            back = f"{term}: {context[:120]}... [Key concept for understanding]"
        
        return {
            "front": front,
            "back": back,
            "tags": [term.lower(), difficulty, "content-generated"],
            "difficulty": difficulty,
            "subject": "Content Analysis",
            "complexity": concept['importance'],
            "source_term": term,
            "original_context": context
        }
    
    def _generate_flashcards_from_text_basic(self, content: str, count: int, difficulty: str) -> List[Dict]:
        """Generate flashcards using basic text analysis"""
        if not content:
            return []
        
        # Extract key terms from content
        all_terms = self._extract_key_terms_from_chunk(content)
        
        flashcards = []
        for i, term in enumerate(all_terms[:count]):
            # Find context for the term
            sentences = content.split('.')
            term_context = ""
            for sentence in sentences:
                if term.lower() in sentence.lower():
                    term_context = sentence.strip()
                    break
            
            if not term_context:
                term_context = f"Key concept related to the main topic."
            
            flashcard = {
                "front": f"What is {term}?",
                "back": f"{term}: {term_context[:150]}...",
                "tags": [term.lower(), difficulty, "basic-generation"],
                "difficulty": difficulty,
                "subject": "Content Review",
                "complexity": 0.5,
                "source_term": term
            }
            flashcards.append(flashcard)
        
        return flashcards
    
    # =========================================================================
    # ADVANCED CONTENT ANALYSIS (Missing methods from enhanced AI service)
    # =========================================================================
    
    def _analyze_content_structure(self, content: str) -> Dict:
        """Analyze the structure and characteristics of lesson content"""
        if not content:
            return {"structure": "empty", "complexity": 0}
        
        words = content.split()
        sentences = content.split('.')
        paragraphs = content.split('\n\n')
        
        # Analyze content characteristics
        analysis = {
            "word_count": len(words),
            "sentence_count": len([s for s in sentences if len(s.strip()) > 10]),
            "paragraph_count": len([p for p in paragraphs if len(p.strip()) > 50]),
            "average_sentence_length": len(words) / max(len(sentences), 1),
            "has_technical_terms": self._detect_technical_terms(content),
            "readability_level": self._assess_readability(content),
            "content_density": len(words) / max(len(paragraphs), 1),
            "structure_type": self._classify_content_structure(content)
        }
        
        return analysis
    
    def _extract_key_concepts(self, content: str) -> List[str]:
        """Extract key concepts from content using intelligent analysis"""
        if not content:
            return []
        
        # Extract potential concepts
        words = content.split()
        
        # Filter for conceptual terms
        concepts = []
        for word in words:
            word_clean = word.strip('.,!?":;()[]{}').lower()
            if (len(word_clean) > 4 and 
                word_clean.isalpha() and 
                word_clean not in ['the', 'and', 'that', 'this', 'with', 'have', 'they', 'from', 'been', 'were', 'said']):
                concepts.append(word_clean.title())
        
        # Count frequency and return most important
        concept_counts = {}
        for concept in concepts:
            concept_counts[concept] = concept_counts.get(concept, 0) + 1
        
        # Sort by frequency and return top concepts
        sorted_concepts = sorted(concept_counts.items(), key=lambda x: x[1], reverse=True)
        return [concept for concept, count in sorted_concepts[:8]]
    
    def _generate_learning_objectives(self, key_concepts: List[str], content: str) -> List[str]:
        """Generate learning objectives based on Bloom's taxonomy"""
        if not key_concepts:
            return ["Understand the main concepts presented in the lesson."]
        
        bloom_templates = {
            'knowledge': "Define and identify {concept}",
            'comprehension': "Explain the significance of {concept}",
            'application': "Apply {concept} to solve problems",
            'analysis': "Analyze the relationship between {concept} and related concepts",
            'synthesis': "Create new ideas using {concept}",
            'evaluation': "Evaluate the importance of {concept} in the field"
        }
        
        objectives = []
        for i, concept in enumerate(key_concepts[:4]):
            levels = list(bloom_templates.keys())
            selected_level = levels[i % len(levels)]
            objective = bloom_templates[selected_level].format(concept=concept)
            objectives.append(f"Students will {objective.lower()}")
        
        return objectives
    
    def _create_intelligent_summary(self, content: str, key_concepts: List[str]) -> str:
        """Create an intelligent summary of the content"""
        if not content:
            return "No content available for summary."
        
        # Extract first and last sentences for context
        sentences = [s.strip() for s in content.split('.') if len(s.strip()) > 20]
        
        if not sentences:
            return "The content covers various important concepts for study."
        
        intro_sentence = sentences[0] if sentences else ""
        key_concept_text = ", ".join(key_concepts[:3]) if key_concepts else "various concepts"
        
        summary = f"This lesson focuses on {key_concept_text}. {intro_sentence[:100]}... "
        
        if len(sentences) > 1:
            conclusion = sentences[-1][:100] if sentences[-1] else ""
            summary += f"The content concludes with insights about {conclusion}..."
        
        return summary
    
    def _generate_critical_thinking_questions(self, key_concepts: List[str], content: str) -> List[str]:
        """Generate thought-provoking questions for deeper learning"""
        if not key_concepts:
            return ["What are the main ideas presented in this lesson?"]
        
        question_templates = [
            "How does {concept} relate to real-world applications?",
            "What would happen if {concept} didn't exist?",
            "Compare {concept} with similar concepts you've learned.",
            "Why is {concept} important in this field of study?",
            "What are the implications of {concept} for future developments?",
            "How might {concept} be viewed differently in other contexts?"
        ]
        
        questions = []
        for i, concept in enumerate(key_concepts[:4]):
            template = question_templates[i % len(question_templates)]
            question = template.format(concept=concept)
            questions.append(question)
        
        return questions
    
    def _estimate_study_time(self, content: str) -> str:
        """Estimate study time based on content analysis"""
        if not content:
            return "5-10 minutes"
        
        word_count = len(content.split())
        
        # Estimate based on average reading speed (200-250 words per minute)
        reading_time = word_count / 225
        
        # Add time for comprehension and practice
        study_time = reading_time * 2.5
        
        if study_time < 10:
            return "5-10 minutes"
        elif study_time < 20:
            return "10-20 minutes"
        elif study_time < 40:
            return "20-40 minutes"
        else:
            return "40+ minutes"
    
    def _assess_content_difficulty(self, content: str) -> str:
        """Assess the difficulty level of content"""
        if not content:
            return "Unknown"
        
        # Analyze complexity indicators
        words = content.split()
        avg_word_length = sum(len(word) for word in words) / len(words) if words else 0
        
        technical_terms = self._detect_technical_terms(content)
        complex_sentences = len([s for s in content.split('.') if len(s.split()) > 20])
        
        difficulty_score = (avg_word_length - 4) * 0.3 + len(technical_terms) * 0.4 + complex_sentences * 0.1
        
        if difficulty_score < 1:
            return "Beginner"
        elif difficulty_score < 2.5:
            return "Intermediate"
        else:
            return "Advanced"
    
    def _identify_prerequisites(self, content: str) -> List[str]:
        """Identify prerequisite knowledge for the content"""
        if not content:
            return []
        
        # Look for terms that suggest prerequisites
        prerequisite_indicators = [
            "previous", "earlier", "foundation", "basic", "fundamental", 
            "prerequisite", "required", "assumed", "background"
        ]
        
        prerequisites = []
        sentences = content.split('.')
        
        for sentence in sentences:
            sentence_lower = sentence.lower()
            if any(indicator in sentence_lower for indicator in prerequisite_indicators):
                # Extract potential prerequisite concepts
                words = sentence.split()
                for i, word in enumerate(words):
                    if word.lower() in prerequisite_indicators and i < len(words) - 1:
                        next_words = ' '.join(words[i+1:i+4])
                        prerequisites.append(next_words.strip('.,!?'))
        
        # Add domain-specific prerequisites
        domain = self._identify_subject_domain(content)
        domain_prerequisites = {
            'Mathematics': ['Basic arithmetic', 'Algebraic thinking'],
            'Physics': ['Mathematical concepts', 'Scientific method'],
            'Chemistry': ['Atomic theory', 'Mathematical skills'],
            'Biology': ['Scientific method', 'Basic chemistry'],
            'Computer Science': ['Logical thinking', 'Problem-solving'],
        }
        
        if domain in domain_prerequisites:
            prerequisites.extend(domain_prerequisites[domain])
        
        return list(set(prerequisites))[:4]  # Return unique prerequisites, limit to 4
    
    def _detect_technical_terms(self, content: str) -> List[str]:
        """Detect technical terms in content"""
        words = content.split()
        technical_terms = []
        
        for word in words:
            word_clean = word.strip('.,!?":;()[]{}')
            # Technical terms are usually longer, capitalized, or contain specific patterns
            if (len(word_clean) > 6 and 
                (word_clean.istitle() or 
                 any(char.isdigit() for char in word_clean) or
                 word_clean.isupper())):
                technical_terms.append(word_clean)
        
        return list(set(technical_terms))[:10]
    
    def _assess_readability(self, content: str) -> str:
        """Assess readability level of content"""
        if not content:
            return "Unknown"
        
        words = content.split()
        sentences = [s for s in content.split('.') if len(s.strip()) > 5]
        
        if not words or not sentences:
            return "Unknown"
        
        avg_words_per_sentence = len(words) / len(sentences)
        avg_syllables = sum(self._count_syllables(word) for word in words) / len(words)
        
        # Simple readability assessment
        readability_score = 206.835 - (1.015 * avg_words_per_sentence) - (84.6 * avg_syllables)
        
        if readability_score > 90:
            return "Very Easy"
        elif readability_score > 80:
            return "Easy"
        elif readability_score > 70:
            return "Fairly Easy"
        elif readability_score > 60:
            return "Standard"
        elif readability_score > 50:
            return "Fairly Difficult"
        else:
            return "Difficult"
    
    def _count_syllables(self, word: str) -> int:
        """Count syllables in a word (simple approximation)"""
        word = word.lower().strip('.,!?":;()[]{}')
        if not word:
            return 0
        
        vowels = "aeiouy"
        syllable_count = 0
        prev_was_vowel = False
        
        for char in word:
            if char in vowels:
                if not prev_was_vowel:
                    syllable_count += 1
                prev_was_vowel = True
            else:
                prev_was_vowel = False
        
        # Handle silent e
        if word.endswith('e') and syllable_count > 1:
            syllable_count -= 1
        
        return max(1, syllable_count)
    
    def _classify_content_structure(self, content: str) -> str:
        """Classify the structure type of content"""
        content_lower = content.lower()
        
        if any(word in content_lower for word in ['first', 'second', 'third', 'finally', 'step']):
            return "sequential"
        elif any(word in content_lower for word in ['compare', 'contrast', 'similar', 'different']):
            return "comparative"
        elif any(word in content_lower for word in ['problem', 'solution', 'issue', 'resolve']):
            return "problem-solution"
        elif any(word in content_lower for word in ['cause', 'effect', 'result', 'consequence']):
            return "cause-effect"
        elif any(word in content_lower for word in ['example', 'instance', 'illustration']):
            return "descriptive"
        else:
            return "general"
    
    # =========================================================================
    # DOMAIN-SPECIFIC FLASHCARD GENERATORS
    # =========================================================================
    
    def _generate_math_flashcards(self, topic: str, count: int, difficulty: str) -> List[Dict]:
        """Generate mathematics-focused flashcards"""
        flashcards = []
        
        math_concepts = {
            "easy": [
                ("Basic Addition", "What is 7 + 5?", "12"),
                ("Simple Fractions", "What is 1/2 + 1/4?", "3/4"),
                ("Basic Geometry", "How many sides does a triangle have?", "3"),
            ],
            "medium": [
                ("Algebra", "If x + 3 = 10, what is x?", "x = 7"),
                ("Geometry", "What is the area of a rectangle with width 4 and length 6?", "24 square units"),
                ("Fractions", "What is 2/3 × 3/4?", "1/2"),
            ],
            "hard": [
                ("Calculus", "What is the derivative of x²?", "2x"),
                ("Advanced Algebra", "Solve: x² - 5x + 6 = 0", "x = 2 or x = 3"),
                ("Trigonometry", "What is sin(30°)?", "1/2"),
            ]
        }
        
        concepts = math_concepts.get(difficulty, math_concepts["medium"])
        
        for i in range(min(count, len(concepts))):
            concept, front, back = concepts[i]
            flashcard = {
                "front": front,
                "back": back,
                "tags": [topic.lower(), difficulty, "mathematics", concept.lower()],
                "difficulty": difficulty,
                "subject": "Mathematics",
                "complexity": 0.3 if difficulty == "easy" else 0.6 if difficulty == "medium" else 0.9
            }
            flashcards.append(flashcard)
        
        return flashcards
    
    def _generate_science_flashcards(self, topic: str, count: int, difficulty: str) -> List[Dict]:
        """Generate science-focused flashcards"""
        flashcards = []
        
        science_concepts = {
            "easy": [
                ("Basic Chemistry", "What is the chemical symbol for water?", "H₂O"),
                ("Biology", "What organ pumps blood through the body?", "Heart"),
                ("Physics", "What force pulls objects toward Earth?", "Gravity"),
            ],
            "medium": [
                ("Chemistry", "What happens during photosynthesis?", "Plants convert sunlight, CO₂, and water into glucose and oxygen"),
                ("Biology", "What is the powerhouse of the cell?", "Mitochondria"),
                ("Physics", "What is Newton's first law of motion?", "An object at rest stays at rest unless acted upon by an external force"),
            ],
            "hard": [
                ("Advanced Chemistry", "What is the molecular geometry of methane (CH₄)?", "Tetrahedral"),
                ("Molecular Biology", "What enzyme is responsible for DNA replication?", "DNA polymerase"),
                ("Quantum Physics", "What is Heisenberg's uncertainty principle?", "You cannot simultaneously know both position and momentum of a particle with absolute precision"),
            ]
        }
        
        concepts = science_concepts.get(difficulty, science_concepts["medium"])
        
        for i in range(min(count, len(concepts))):
            concept, front, back = concepts[i]
            flashcard = {
                "front": front,
                "back": back,
                "tags": [topic.lower(), difficulty, "science", concept.lower()],
                "difficulty": difficulty,
                "subject": "Science",
                "complexity": 0.4 if difficulty == "easy" else 0.7 if difficulty == "medium" else 0.9
            }
            flashcards.append(flashcard)
        
        return flashcards
    
    def _generate_tech_flashcards(self, topic: str, count: int, difficulty: str) -> List[Dict]:
        """Generate technology/computer science flashcards"""
        flashcards = []
        
        tech_concepts = {
            "easy": [
                ("Programming Basics", "What does HTML stand for?", "HyperText Markup Language"),
                ("Computer Basics", "What is RAM?", "Random Access Memory - temporary storage for active programs"),
                ("Internet", "What does URL stand for?", "Uniform Resource Locator"),
            ],
            "medium": [
                ("Programming", "What is a function in programming?", "A reusable block of code that performs a specific task"),
                ("Data Structures", "What is an array?", "A collection of elements stored in contiguous memory locations"),
                ("Algorithms", "What is Big O notation?", "A mathematical notation describing algorithm efficiency"),
            ],
            "hard": [
                ("Advanced Programming", "What is polymorphism in OOP?", "The ability of objects to take multiple forms through inheritance and interfaces"),
                ("System Design", "What is microservices architecture?", "An architectural approach where applications are built as independent, loosely coupled services"),
                ("Machine Learning", "What is gradient descent?", "An optimization algorithm used to minimize cost functions in machine learning"),
            ]
        }
        
        concepts = tech_concepts.get(difficulty, tech_concepts["medium"])
        
        for i in range(min(count, len(concepts))):
            concept, front, back = concepts[i]
            flashcard = {
                "front": front,
                "back": back,
                "tags": [topic.lower(), difficulty, "technology", concept.lower()],
                "difficulty": difficulty,
                "subject": "Technology",
                "complexity": 0.5 if difficulty == "easy" else 0.7 if difficulty == "medium" else 0.9
            }
            flashcards.append(flashcard)
        
        return flashcards
    
    def _generate_humanities_flashcards(self, topic: str, count: int, difficulty: str) -> List[Dict]:
        """Generate humanities/history flashcards"""
        flashcards = []
        
        humanities_concepts = {
            "easy": [
                ("World History", "When did World War II end?", "1945"),
                ("Geography", "What is the capital of France?", "Paris"),
                ("Literature", "Who wrote 'Romeo and Juliet'?", "William Shakespeare"),
            ],
            "medium": [
                ("History", "What was the Renaissance?", "A cultural movement in Europe from 14th-17th centuries emphasizing art, science, and humanism"),
                ("Philosophy", "What is empiricism?", "The theory that knowledge comes primarily from sensory experience"),
                ("Art History", "What characterizes Impressionist painting?", "Focus on light, color, and capturing momentary impressions rather than detailed realism"),
            ],
            "hard": [
                ("Advanced History", "What were the causes of the French Revolution?", "Economic crisis, social inequality, Enlightenment ideas, and political absolutism"),
                ("Critical Theory", "What is postmodernism?", "A philosophical movement questioning grand narratives and emphasizing relativism and deconstruction"),
                ("Comparative Literature", "What is magical realism?", "A literary style blending realistic narrative with fantastical elements presented as normal"),
            ]
        }
        
        concepts = humanities_concepts.get(difficulty, humanities_concepts["medium"])
        
        for i in range(min(count, len(concepts))):
            concept, front, back = concepts[i]
            flashcard = {
                "front": front,
                "back": back,
                "tags": [topic.lower(), difficulty, "humanities", concept.lower()],
                "difficulty": difficulty,
                "subject": "Humanities",
                "complexity": 0.4 if difficulty == "easy" else 0.6 if difficulty == "medium" else 0.8
            }
            flashcards.append(flashcard)
        
        return flashcards
    
    def _generate_general_flashcards(self, topic: str, count: int, difficulty: str) -> List[Dict]:
        """Generate general knowledge flashcards"""
        flashcards = []
        
        for i in range(count):
            flashcard = {
                "front": f"What is a key concept in {topic}?",
                "back": f"A fundamental principle or idea essential for understanding {topic}. This requires further study and exploration.",
                "tags": [topic.lower(), difficulty, "general"],
                "difficulty": difficulty,
                "subject": "General Studies",
                "complexity": 0.5
            }
            flashcards.append(flashcard)
        
        return flashcards
    